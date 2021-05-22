import ezagui
import sys

sys.dont_write_bytecode = True
import os
import shutil
import docx
from docx import Document
from docx.shared import Inches, Cm
from docx.oxml.shared import OxmlElement
from docx.enum.text import WD_ALIGN_PARAGRAPH
from docx.enum.table import WD_ALIGN_VERTICAL
from docx.shared import RGBColor
from docx.shared import Pt
from pptx import Presentation
from pptx.oxml.ns import _nsmap
from pptx.shapes.group import GroupShape
from pptx.enum.shapes import MSO_SHAPE_TYPE
from pptx.enum.shapes import PP_PLACEHOLDER


pptxfldr = ""

_nsmap["adec"] = "http://schemas.microsoft.com/office/drawing/2017/decorative"
document = Document()
document.add_heading("Accessibility Report", 0)


def create_report_doc(pptxfldr):
    pptxfldr = pptxfldr
    sections = document.sections
    for section in sections:
        section.top_margin = Cm(0.5)
        section.bottom_margin = Cm(0.5)
        section.left_margin = Cm(1.2)
        section.right_margin = Cm(1)


n = 0
d = 0
img_names = []
sldnums = []
img_paths = []
group_img_paths = ""
grptxt = ""
group_text = []
auto_desc_txts = []


def extract_group_txt(groupshape, grptxt):
    grptxt = ""
    for shape in groupshape:
        if shape.shape_type == MSO_SHAPE_TYPE.GROUP:
            grptxt = extract_group_txt(shape.shapes, grptxt)
        else:
            if hasattr(shape, "text"):
                grptxt += shape.text + " "
            else:
                grptxt += ""
    return grptxt


def extract_group_img(groupshape, group_img_paths, pptxfldr):
    global d
    group_img_paths = ""
    for shape in groupshape:
        if shape.shape_type == MSO_SHAPE_TYPE.GROUP:
            group_img_paths = extract_group_img(shape.shapes, group_img_paths, pptxfldr)
        else:
            if hasattr(shape, "image"):
                image = shape.image
                image_bytes = image.blob
                image_filename = "group{:04d}.{}".format(d, image.ext)
                group_img_paths += image_filename + " "
                d += 1
                with open(pptxfldr + "/tmp/" + image_filename, "wb") as f:
                    f.write(image_bytes)
            else:
                group_img_paths += ""
    return group_img_paths


def extract_shape(shape, pptxfldr):
    global n
    global grptxt
    global group_img_paths
    if shape.shape_type == MSO_SHAPE_TYPE.GROUP:
        group_img_paths = extract_group_img(shape.shapes, group_img_paths, pptxfldr)
        grptxt = extract_group_txt(shape.shapes, grptxt)
        group_text.append(grptxt)
        img_paths.append(group_img_paths)

    elif shape.shape_type == MSO_SHAPE_TYPE.PICTURE:
        group_text.append("")
        image = shape.image
        image_bytes = image.blob
        image_filename = "image{:04d}.{}".format(n, image.ext)
        img_paths.append(image_filename)
        with open(pptxfldr + "/tmp/" + image_filename, "wb") as f:
            f.write(image_bytes)
    else:
        group_text.append("")
        img_paths.append("")
    
    n += 1 


# Function obsolete, implement into extract_shape
def valid_shape_type(shape):
    if shape.shape_type == MSO_SHAPE_TYPE.PICTURE:
        if ".jpg" or ".png" in picture.name:
            return True
        else:
            return False
    elif shape.shape_type == MSO_SHAPE_TYPE.GROUP:
        return True
    elif shape.shape_type == MSO_SHAPE_TYPE.CHART:
        return True
    elif shape.shape_type == MSO_SHAPE_TYPE.DIAGRAM:
        return True
    elif shape.shape_type == MSO_SHAPE_TYPE.TABLE:
        return True
    elif shape.shape_type == MSO_SHAPE_TYPE.MEDIA:
        return True
    elif shape.shape_type == MSO_SHAPE_TYPE.WEB_VIDEO:
        return True
    else:
        return False


def is_decorative(shape):
    cNvPr = shape._element._nvXxPr.cNvPr
    adec_decoratives = cNvPr.xpath(".//adec:decorative[@val='1']")
    if adec_decoratives:
        return True


# Remove unnecessary elses (if applicable), check if you can use None instead of 0
def has_alt_text(shape):
    autoDesc = "Description automatically generated"
    if shape.alt_text:
        if autoDesc in shape.alt_text:
            auto_desc_txts.append(shape.alt_text)
            return False
        else:
            return True
    else:
        auto_desc_txts.append("")


def keep_table_on_one_page(document):
    """https://github.com/python-openxml/python-docx/issues/245#event-621236139
    Globally prevent table cells from splitting across pages.
    """
    tags = document.element.xpath("//w:tr")
    rows = len(tags)
    for row in range(rows):
        tag = tags[row]  # Specify which <w:r> tag you want
        child = OxmlElement("w:cantSplit")  # Create arbitrary tag
        tag.append(child)  # Append in the new tag


def create_docx_table(filename, pptxfldr):
    f = document.add_paragraph()
    d = f.add_run("File: ")
    r = f.add_run(filename)
    r.bold = True
    font = r.font
    font.name = "Malgun Gothic"
    font.size = Pt(19)
    font = d.font
    font.name = "Malgun Gothic"
    font.size = Pt(19)
    f.alignment = WD_ALIGN_PARAGRAPH.LEFT
    table = document.add_table(rows=1, cols=4, style="Light Grid Accent 1")
    hdr_cells = table.rows[0].cells
    p = hdr_cells[0].paragraphs[0]
    r = p.add_run(style=None)
    font = r.font
    font.name = "Malgun Gothic"
    font.size = Pt(11)
    r.add_text("Figure Name")
    p.alignment = WD_ALIGN_PARAGRAPH.CENTER
    p = hdr_cells[1].paragraphs[0]
    r = p.add_run(style=None)
    font = r.font
    font.name = "Malgun Gothic"
    font.size = Pt(10)
    r.add_text("Slide Number")
    p.alignment = WD_ALIGN_PARAGRAPH.CENTER
    p = hdr_cells[2].paragraphs[0]
    r = p.add_run(style=None)
    font = r.font
    font.name = "Malgun Gothic"
    font.size = Pt(11)
    r.add_text("Figure Image/Text")
    p.alignment = WD_ALIGN_PARAGRAPH.CENTER
    p = hdr_cells[3].paragraphs[0]
    r = p.add_run(style=None)
    font = r.font
    font.name = "Malgun Gothic"
    font.size = Pt(11)
    r.add_text("Alt-Text")
    p.alignment = WD_ALIGN_PARAGRAPH.CENTER

    for (img_name, sldnum, img_path, grouptext, auto_desc_txt) in zip(
        img_names, sldnums, img_paths, group_text, auto_desc_txts):
        row_cells = table.add_row().cells
        p = row_cells[0].paragraphs[0]
        r = p.add_run(style=None)
        r.add_text(img_name)
        font = r.font
        font.name = "Malgun Gothic"
        font.size = Pt(12.5)
        p.alignment = WD_ALIGN_PARAGRAPH.CENTER
        row_cells[0].vertical_alignment = WD_ALIGN_VERTICAL.CENTER
        p = row_cells[1].paragraphs[0]
        r = p.add_run(style=None)
        r.add_text(str(sldnum))
        font = r.font
        font.name = "Malgun Gothic"
        font.size = Pt(16)
        p.alignment = WD_ALIGN_PARAGRAPH.CENTER
        row_cells[1].vertical_alignment = WD_ALIGN_VERTICAL.CENTER
        p = row_cells[2].paragraphs[0]
        r = p.add_run(style=None)
        sngle_path = img_path.split(" ")
        for path in sngle_path:
            if path != "":
                if ".wmf" not in path:
                    print(path)
                    r.add_picture(pptxfldr + "/tmp/" + path, width=Inches(2))
                    p.alignment = WD_ALIGN_PARAGRAPH.CENTER
            else:
                r.add_break()
                r.add_text("Figure type not supported")
                r.add_break()
                r.add_text("Please insert manually!")
                r.add_break()
                r.bold = True
                font = r.font
                font.color.rgb = RGBColor(255, 0, 0)
                font.size = Pt(13)
                p.alignment = WD_ALIGN_PARAGRAPH.CENTER
        if grouptext != "":
            r.add_break()
            font = r.font
            font.name = "Malgun Gothic"
            font.size = Pt(9)
            r.add_text(grouptext)
            p.alignment = WD_ALIGN_PARAGRAPH.CENTER
            row_cells[2].vertical_alignment = WD_ALIGN_VERTICAL.CENTER

        p = row_cells[3].paragraphs[0]
        r = p.add_run(style=None)
        font = r.font
        font.name = "Malgun Gothic"
        font.size = Pt(8)
        r.add_text(auto_desc_txt)
        p.alignment = WD_ALIGN_PARAGRAPH.CENTER
        row_cells[3].vertical_alignment = WD_ALIGN_VERTICAL.CENTER

        for cells in row_cells:
            p = cells.paragraphs[0]
            p.style = None

    widths = (Inches(1.3), Inches(1.1), Inches(2.9), Inches(2.4))
    for row in table.rows:
        for idx, width in enumerate(widths):
            row.cells[idx].width = width

    document.add_page_break()


def check_file_accessibility(prs, pptxfldr):
    for slide in prs.slides:
        for shape in slide.shapes:
            if (
                valid_shape_type(shape)
                and not is_decorative(shape)
                and not has_alt_text(shape)
            ):
                extract_shape(shape, pptxfldr)
                yield shape
                img_names.append(shape.name)
                sldnum = prs.slides.index(slide) + 1
                sldnums.append(str(sldnum))
                n = 0


def iter_files(pptxfldr, update_status):
    for filename in os.listdir(pptxfldr):
        if filename.endswith(".pptx"):
            filename = filename.replace("~$", "")
            update_status(filename)
            for picture in check_file_accessibility(
                Presentation(pptxfldr + "/" + filename), pptxfldr
            ):
                continue
            pptxfldr = pptxfldr
            print(img_names)
            print(sldnums)
            print(img_paths)
            print(img_names)
            create_docx_table(filename, pptxfldr)
            img_names.clear()
            sldnums.clear()
            img_paths.clear()
            grptxt = ""
            group_text.clear()
            auto_desc_txts.clear()
            print(img_names)
            print(sldnums)
            print(img_paths)
            print(img_names)


def save_report(pptxfldr, savepath):
    keep_table_on_one_page(document)
    if savepath:
        document.save(savepath)
    else:
        document.save(pptxfldr + "/sample-report.docx")
